#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
🎨 PPT AI 重新設計工具 - Streamlit 互動應用
Q3: PPT 換版型（AI 重新設計）

功能:
  • 上傳 PPT 檔案
  • 選擇設計風格 (現代科技風 / 商務沉穩風)
  • 即時預覽配色方案
  • 下載設計後的 PPT

作者: AI Assistant
日期: 2025/12/24
"""

import streamlit as st
import io
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import tempfile
import os

# ============================================================================
# 頁面配置
# ============================================================================

st.set_page_config(
    page_title="🎨 PPT AI 重新設計工具",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 自訂 CSS
st.markdown("""
<style>
    .main {
        padding-top: 2rem;
    }
    .design-box {
        border: 2px solid #ddd;
        border-radius: 8px;
        padding: 20px;
        margin: 10px 0;
    }
    .modern-style {
        background: linear-gradient(135deg, #141e3c 0%, #9664ff 100%);
        color: white;
    }
    .business-style {
        background: linear-gradient(135deg, #f0f0f0 0%, #e8e8e8 100%);
        color: #3c3c3c;
    }
    .success-box {
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        color: #155724;
        padding: 12px;
        border-radius: 4px;
        margin: 10px 0;
    }
    .info-box {
        background-color: #d1ecf1;
        border: 1px solid #bee5eb;
        color: #0c5460;
        padding: 12px;
        border-radius: 4px;
        margin: 10px 0;
    }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# 顏色方案定義
# ============================================================================

STYLE_CONFIGS = {
    "現代科技風": {
        "name": "Modern Tech",
        "description": "深藍背景 + 亮藍紫色強調 | 高對比、科技感強",
        "colors": {
            "bg": (20, 30, 60),          # 深藍
            "text": (255, 255, 255),     # 白色
            "accent1": (100, 200, 255),  # 亮藍
            "accent2": (150, 100, 255),  # 紫色
        },
        "preview_bg": "#141E3C",
        "preview_text": "#FFFFFF",
        "emoji": "🚀"
    },
    "商務沉穩風": {
        "name": "Business Professional",
        "description": "淺灰背景 + 金色標題 | 專業可信、企業級",
        "colors": {
            "bg": (240, 240, 240),       # 淺灰
            "text": (60, 60, 60),        # 深灰文字
            "accent1": (184, 134, 11),   # 金色
            "accent2": (45, 45, 45),     # 深灰
        },
        "preview_bg": "#F0F0F0",
        "preview_text": "#3C3C3C",
        "emoji": "💼"
    }
}

# ============================================================================
# 設計功能函數
# ============================================================================

def apply_style_to_ppt(input_ppt_bytes, style_name):
    """
    套用設計風格到 PPT
    
    Args:
        input_ppt_bytes: PPT 檔案的 bytes
        style_name: 風格名稱 ('現代科技風' 或 '商務沉穩風')
    
    Returns:
        bytes: 設計後的 PPT 檔案
    """
    # 讀取 PPT
    prs = Presentation(io.BytesIO(input_ppt_bytes))
    style_config = STYLE_CONFIGS[style_name]
    colors = style_config["colors"]
    
    # 套用風格到每一張投影片
    for i, slide in enumerate(prs.slides):
        # 設置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors["bg"])
        
        # 調整文字樣式
        for j, shape in enumerate(slide.shapes):
            try:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 商務風格: 第一個形狀用金色
                            if style_name == "商務沉穩風" and j == 0:
                                run.font.color.rgb = RGBColor(*colors["accent1"])
                                run.font.bold = True
                            else:
                                run.font.color.rgb = RGBColor(*colors["text"])
            except Exception as e:
                pass  # 跳過無法處理的形狀
    
    # 保存到 bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

def get_file_info(ppt_bytes):
    """獲取 PPT 檔案資訊"""
    try:
        prs = Presentation(io.BytesIO(ppt_bytes))
        return {
            "slides": len(prs.slides),
            "size_mb": len(ppt_bytes) / (1024 * 1024),
            "width": prs.slide_width.inches,
            "height": prs.slide_height.inches
        }
    except Exception as e:
        return None

# ============================================================================
# 主要應用
# ============================================================================

def main():
    # 標題
    st.markdown("""
    <div style="text-align: center;">
        <h1>🎨 PPT AI 重新設計工具</h1>
        <p style="font-size: 18px; color: #666;">
            使用 AI 快速為您的 PowerPoint 進行版型重新設計
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # 側邊欄
    with st.sidebar:
        st.markdown("### 📋 使用說明")
        st.info("""
        **步驟:**
        1. 上傳您的 PowerPoint 檔案 (.pptx)
        2. 選擇喜歡的設計風格
        3. 點擊「應用設計」按鈕
        4. 下載重新設計的 PPT
        
        **支援格式:** .pptx (PowerPoint 2007+)
        """)
        
        st.markdown("---")
        st.markdown("### 🎯 風格特點")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("""
            **🚀 現代科技風**
            - 深藍 + 白文字
            - 高對比、科技感
            - 適合: 技術展示
            """)
        
        with col2:
            st.markdown("""
            **💼 商務沉穩風**
            - 淺灰 + 金色標題
            - 專業可信
            - 適合: 商業洽談
            """)
    
    # 主要內容
    col1, col2 = st.columns([2, 1], gap="large")
    
    with col1:
        st.markdown("### 📤 上傳 PowerPoint 檔案")
        uploaded_file = st.file_uploader(
            "選擇 .pptx 檔案",
            type=["pptx"],
            help="支援 Microsoft PowerPoint 2007 及以上版本"
        )
        
        if uploaded_file is not None:
            # 顯示檔案資訊
            file_info = get_file_info(uploaded_file.getvalue())
            
            if file_info:
                st.success("✅ 檔案讀取成功")
                
                col_info1, col_info2, col_info3 = st.columns(3)
                with col_info1:
                    st.metric("投影片", file_info["slides"])
                with col_info2:
                    st.metric("檔案大小", f"{file_info['size_mb']:.2f} MB")
                with col_info3:
                    st.metric("尺寸", f"{file_info['width']:.1f}\" × {file_info['height']:.1f}\"")
            
            # 風格選擇
            st.markdown("### 🎨 選擇設計風格")
            
            selected_style = st.radio(
                "選擇一個風格",
                list(STYLE_CONFIGS.keys()),
                horizontal=False,
                label_visibility="collapsed"
            )
            
            # 風格預覽
            style_info = STYLE_CONFIGS[selected_style]
            st.markdown(f"""
            <div class="design-box" style="background-color: {style_info['preview_bg']}; color: {style_info['preview_text']}; border: 3px solid {style_info['preview_bg']};">
                <h3>{style_info['emoji']} {selected_style}</h3>
                <p style="margin: 0;">{style_info['description']}</p>
            </div>
            """, unsafe_allow_html=True)
            
            # 應用設計按鈕
            if st.button("✨ 應用設計", use_container_width=True, type="primary"):
                with st.spinner("🔄 正在處理您的 PPT..."):
                    try:
                        # 應用風格
                        new_ppt = apply_style_to_ppt(
                            uploaded_file.getvalue(),
                            selected_style
                        )
                        
                        # 顯示成功訊息
                        st.markdown('<div class="success-box">✅ 設計應用成功！可以下載您的新 PPT</div>', 
                                   unsafe_allow_html=True)
                        
                        # 下載按鈕
                        st.download_button(
                            label="📥 下載設計後的 PPT",
                            data=new_ppt,
                            file_name=f"{selected_style}_{uploaded_file.name}",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                        
                        # 顯示預覽配色
                        st.markdown("#### 🎨 配色方案")
                        colors = style_info["colors"]
                        col_colors = st.columns(len(colors))
                        color_names = ["背景", "文字", "強調1", "強調2"]
                        
                        for i, (name, rgb) in enumerate(colors.items()):
                            with col_colors[i]:
                                hex_color = "#{:02x}{:02x}{:02x}".format(*rgb)
                                st.markdown(f"""
                                <div style="background-color: {hex_color}; padding: 20px; border-radius: 8px; text-align: center;">
                                    <p style="color: {'white' if name == 'bg' else 'black'}; margin: 0; font-weight: bold;">{color_names[i]}</p>
                                    <p style="color: {'white' if name == 'bg' else 'black'}; margin: 5px 0; font-size: 12px;">{hex_color}</p>
                                </div>
                                """, unsafe_allow_html=True)
                    
                    except Exception as e:
                        st.error(f"❌ 處理失敗: {str(e)}")
    
    with col2:
        st.markdown("### ℹ️ 設計資訊")
        
        st.info("""
        **現代科技風** 🚀
        
        RGB 配色:
        - 背景: (20, 30, 60)
        - 文字: (255, 255, 255)
        - 強調: (100, 200, 255)
        
        特點:
        ✓ 科技感強
        ✓ 高對比度
        ✓ 視覺衝擊力大
        """)
        
        st.info("""
        **商務沉穩風** 💼
        
        RGB 配色:
        - 背景: (240, 240, 240)
        - 文字: (60, 60, 60)
        - 標題: (184, 134, 11)
        
        特點:
        ✓ 專業感強
        ✓ 易於列印
        ✓ 企業級配色
        """)
    
    # 頁腳
    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; color: #666; font-size: 12px;">
        <p>🎓 作業 Q3: PPT 換版型（AI 重新設計）</p>
        <p>使用 Python + python-pptx + Streamlit</p>
        <p>© 2025 智慧計算系統課程</p>
    </div>
    """, unsafe_allow_html=True)

# ============================================================================
# 執行應用
# ============================================================================

if __name__ == "__main__":
    main()
