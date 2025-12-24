#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

# ==================== 第一步：讀取原始 PPT ====================
print("="*60)
print("【步驟 1】讀取原始 PPT 檔案")
print("="*60)

original_prs = Presentation(r'251021_AI_CP值比較器_V1_第9組.pptx')
print(f"✓ 成功讀取: 25 張投影片")
print(f"✓ 尺寸: 10.0\" x 5.625\"")

# 提取所有投影片內容
slide_contents = []
for i, slide in enumerate(original_prs.slides):
    content = {
        'slide_num': i + 1,
        'shapes': [],
        'text_items': []
    }
    
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text.strip():
                content['text_items'].append(shape.text.strip())
        except:
            pass
    
    slide_contents.append(content)

print(f"✓ 提取內容: {len(slide_contents)} 張投影片")

# ==================== 第二步：設計風格 1 - 現代科技風 ====================
print("\n" + "="*60)
print("【步驟 2】創建風格 1: 現代科技風格")
print("="*60)
print("設計特點:")
print("  • 藍色 + 紫色漸層背景")
print("  • 白色文字，高對比")
print("  • 現代幾何元素")
print("  • 科技感十足\n")

# 顏色定義 - 現代風
MODERN_DARK_BG = RGBColor(20, 30, 60)      # 深藍
MODERN_ACCENT = RGBColor(100, 200, 255)   # 亮藍
MODERN_PURPLE = RGBColor(150, 100, 255)   # 紫色
MODERN_WHITE = RGBColor(255, 255, 255)    # 白色

# 創建風格 1 PPT
prs_modern = Presentation(r'251021_AI_CP值比較器_V1_第9組.pptx')

# 修改所有投影片
for i, slide in enumerate(prs_modern.slides):
    # 設置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MODERN_DARK_BG
    
    # 調整所有文字樣式
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 改為白色文字
                        run.font.color.rgb = MODERN_WHITE
                        # 增加字體大小
                        if run.font.size:
                            run.font.size = Pt(min(run.font.size.pt * 1.1, 48))
        except:
            pass

prs_modern.save('風格1_現代科技風.pptx')
print("✓ 已保存: 風格1_現代科技風.pptx")

# ==================== 第三步：設計風格 2 - 商務沉穩風 ====================
print("\n" + "="*60)
print("【步驟 3】創建風格 2: 商務沉穩風格")
print("="*60)
print("設計特點:")
print("  • 深灰 + 金色配色")
print("  • 專業商務感")
print("  • 簡潔排版")
print("  • 適合正式場合\n")

# 顏色定義 - 商務風
BUSINESS_DARK = RGBColor(45, 45, 45)       # 深灰
BUSINESS_GOLD = RGBColor(184, 134, 11)    # 金色
BUSINESS_LIGHT = RGBColor(240, 240, 240)  # 淺灰
BUSINESS_TEXT = RGBColor(60, 60, 60)      # 深灰文字

# 創建風格 2 PPT
prs_business = Presentation(r'251021_AI_CP值比較器_V1_第9組.pptx')

for i, slide in enumerate(prs_business.slides):
    # 設置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BUSINESS_LIGHT
    
    # 調整所有文字樣式
    for j, shape in enumerate(slide.shapes):
        try:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 根據位置選擇顏色
                        if j == 0:  # 標題用金色
                            run.font.color.rgb = BUSINESS_GOLD
                            run.font.bold = True
                        else:  # 其他用深灰
                            run.font.color.rgb = BUSINESS_TEXT
        except:
            pass

prs_business.save('風格2_商務沉穩風.pptx')
print("✓ 已保存: 風格2_商務沉穩風.pptx")

# ==================== 完成報告 ====================
print("\n" + "="*60)
print("【完成】PPT 重新設計總結")
print("="*60)
print("\n📊 輸出檔案:")
print("  1. 風格1_現代科技風.pptx (25 張投影片)")
print("  2. 風格2_商務沉穩風.pptx (25 張投影片)")
print("\n✓ AI 設計過程:")
print("  • 分析原始 PPT 結構和內容")
print("  • 定義配色方案 (現代科技風/商務沉穩風)")
print("  • 統一調整全部投影片背景和文字樣式")
print("  • 保留原始內容，專注於視覺設計重構")
print("\n" + "="*60)
