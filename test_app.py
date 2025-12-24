#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Streamlit 應用功能測試
驗證核心設計功能是否正常運作
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

import io
from pptx import Presentation
from pptx.dml.color import RGBColor

def test_ppt_loading():
    """測試 PPT 讀取"""
    print("=" * 60)
    print("【測試 1】PPT 檔案讀取")
    print("=" * 60)
    
    try:
        prs = Presentation('251021_AI_CP值比較器_V1_第9組.pptx')
        print(f"✅ 成功讀取 PPT")
        print(f"   • 投影片數: {len(prs.slides)} 張")
        print(f"   • 尺寸: {prs.slide_width.inches}\" × {prs.slide_height.inches}\"")
        return True
    except Exception as e:
        print(f"❌ 讀取失敗: {str(e)}")
        return False

def test_modern_style():
    """測試現代科技風設計"""
    print("\n" + "=" * 60)
    print("【測試 2】現代科技風設計套用")
    print("=" * 60)
    
    try:
        # 讀取原始檔案
        prs = Presentation('251021_AI_CP值比較器_V1_第9組.pptx')
        
        # 套用現代風格
        modern_bg = RGBColor(20, 30, 60)
        modern_text = RGBColor(255, 255, 255)
        
        for slide in prs.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = modern_bg
            
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = modern_text
                except:
                    pass
        
        # 保存測試檔案
        test_output = io.BytesIO()
        prs.save(test_output)
        test_output.seek(0)
        
        print(f"✅ 現代科技風設計成功套用")
        print(f"   • 背景色: RGB(20, 30, 60)")
        print(f"   • 文字色: RGB(255, 255, 255)")
        print(f"   • 輸出大小: {len(test_output.getvalue()) / (1024*1024):.2f} MB")
        return True
    except Exception as e:
        print(f"❌ 套用失敗: {str(e)}")
        return False

def test_business_style():
    """測試商務沉穩風設計"""
    print("\n" + "=" * 60)
    print("【測試 3】商務沉穩風設計套用")
    print("=" * 60)
    
    try:
        # 讀取原始檔案
        prs = Presentation('251021_AI_CP值比較器_V1_第9組.pptx')
        
        # 套用商務風格
        business_bg = RGBColor(240, 240, 240)
        business_text = RGBColor(60, 60, 60)
        business_accent = RGBColor(184, 134, 11)
        
        for slide in prs.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = business_bg
            
            for j, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if j == 0:
                                    run.font.color.rgb = business_accent
                                    run.font.bold = True
                                else:
                                    run.font.color.rgb = business_text
                except:
                    pass
        
        # 保存測試檔案
        test_output = io.BytesIO()
        prs.save(test_output)
        test_output.seek(0)
        
        print(f"✅ 商務沉穩風設計成功套用")
        print(f"   • 背景色: RGB(240, 240, 240)")
        print(f"   • 文字色: RGB(60, 60, 60)")
        print(f"   • 強調色: RGB(184, 134, 11)")
        print(f"   • 輸出大小: {len(test_output.getvalue()) / (1024*1024):.2f} MB")
        return True
    except Exception as e:
        print(f"❌ 套用失敗: {str(e)}")
        return False

def test_streamlit_import():
    """測試 Streamlit 套件"""
    print("\n" + "=" * 60)
    print("【測試 4】Streamlit 套件檢查")
    print("=" * 60)
    
    try:
        import streamlit as st
        print(f"✅ Streamlit 已安裝")
        print(f"   • 版本: {st.__version__}")
        return True
    except ImportError:
        print(f"⚠️ Streamlit 未安裝")
        print(f"   執行: pip install streamlit")
        return False

def main():
    print("\n")
    print("🔍 Streamlit 應用功能測試")
    print("=" * 60)
    
    results = []
    
    # 運行所有測試
    results.append(("PPT 讀取", test_ppt_loading()))
    results.append(("現代科技風", test_modern_style()))
    results.append(("商務沉穩風", test_business_style()))
    results.append(("Streamlit 套件", test_streamlit_import()))
    
    # 總結
    print("\n" + "=" * 60)
    print("【測試總結】")
    print("=" * 60)
    
    passed = sum(1 for _, result in results if result)
    total = len(results)
    
    for test_name, result in results:
        status = "✅ 通過" if result else "❌ 失敗"
        print(f"{status}: {test_name}")
    
    print(f"\n結果: {passed}/{total} 測試通過")
    
    if passed == total:
        print("\n🎉 所有測試通過！應用可以部署")
        print("\n執行以下命令啟動應用:")
        print("  streamlit run streamlit_app.py")
        return 0
    else:
        print("\n⚠️ 有測試未通過，請檢查依賴")
        return 1

if __name__ == "__main__":
    sys.exit(main())
